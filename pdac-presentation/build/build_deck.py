"""Build the LLM Pancreatic Oncology Clinical Trial System seminar deck.

Produces a 16:9 landscape ``.pptx`` of twenty-three slides:

    1       title slide
    2 - 21  one slide per deposited paper, in chronological order, each
            carrying that paper's cover image and its seven-line evaluation
    22 - 23 references, with every DOI hyperlinked to its resolver

Design constraints this module enforces mechanically rather than by eye:

* The palette is exactly #002f5f, #666666, #407cb9, #e3eaee and #dcdcdc.
* Every slide ground is white or near-white, because all twenty cover images
  are white-page scans and would be framed badly by a dark slide.
* Cover images run four slides on the left, then one on the right, repeating,
  so papers 5, 10, 15 and 20 carry right-hand covers.  ``slide_content.side_for``
  owns that rule.
* Every bullet and numbered item is a single line.  ``assert_single_line``
  measures each one with the metric-compatible Liberation face before the
  deck is written and fails the build if any item would wrap, so a copy edit
  can never silently produce a two-line bullet.
* Page numbers sit in the lower right corner and the author and date sit in
  the footer of every slide.

Run from the repository root::

    python pdac-presentation/build/build_deck.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

import slide_content as content  # noqa: E402

REPO_ROOT = Path(__file__).resolve().parents[2]
PRESENTATION_DIR = REPO_ROOT / "pdac-presentation"
COVER_DIR = PRESENTATION_DIR / "cover-images" / "extracted"
OUTPUT_PPTX = PRESENTATION_DIR / "slides" / "LLM-Pancreatic-Oncology-Clinical-Trial-System.pptx"

# ---------------------------------------------------------------------------
# Palette.  These five values are the whole colour vocabulary of the deck.
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x00, 0x2F, 0x5F)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x40, 0x7C, 0xB9)
PALE = RGBColor(0xE3, 0xEA, 0xEE)
SILVER = RGBColor(0xDC, 0xDC, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x1F, 0x1F)

# Arial and Times New Roman are used because their metrics are the ones the
# fit check measures (Liberation Sans and Liberation Serif are metric clones),
# so the PowerPoint rendering and the LibreOffice PDF rendering agree.
SANS = "Arial"
SERIF = "Times New Roman"

FONT_DIR = Path("/usr/share/fonts/truetype/liberation")
METRIC_FONTS = {
    (SANS, False): FONT_DIR / "LiberationSans-Regular.ttf",
    (SANS, True): FONT_DIR / "LiberationSans-Bold.ttf",
    (SERIF, False): FONT_DIR / "LiberationSerif-Regular.ttf",
    (SERIF, True): FONT_DIR / "LiberationSerif-Bold.ttf",
}
MEASURE_SCALE = 10  # measure at 10x the point size, then divide, for precision

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
CONTENT_RIGHT = Inches(12.833)

HEADER_Y = Inches(0.16)
RULE_Y = Inches(0.42)
HEADLINE_Y = Inches(0.52)
HEADLINE_H = Inches(0.46)
BODY_TOP = Inches(1.16)
BODY_BOTTOM = Inches(6.68)
FOOTER_RULE_Y = Inches(6.86)
FOOTER_TEXT_Y = Inches(6.93)
FOOTER_H = Inches(0.30)

IMAGE_BOX_W = Inches(3.45)
IMAGE_BOX_H = Inches(5.36)
COLUMN_GAP = Inches(0.28)
TEXT_COL_W = Emu(int(CONTENT_RIGHT) - int(MARGIN) - int(IMAGE_BOX_W) - int(COLUMN_GAP))

BODY_PT = 12.5
LABEL_PT = 12.5
BULLET_INDENT = Inches(0.24)
TEXT_INSET = Inches(0.06)
LIST_INSET_L = Inches(0.14)
LIST_INSET_R = Inches(0.06)
FIT_SAFETY_PT = 2.0

BULLET_CHAR = "▪"


# ---------------------------------------------------------------------------
# Measurement
# ---------------------------------------------------------------------------
_font_cache: dict[tuple[str, bool, int], ImageFont.FreeTypeFont] = {}


def _metric_font(typeface: str, bold: bool, size_pt: float) -> ImageFont.FreeTypeFont:
    key = (typeface, bold, int(round(size_pt * MEASURE_SCALE)))
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(str(METRIC_FONTS[(typeface, bold)]), key[2])
    return _font_cache[key]


def text_width_pt(text: str, typeface: str, bold: bool, size_pt: float) -> float:
    """Width of ``text`` in points, measured on the metric-compatible face."""
    return _metric_font(typeface, bold, size_pt).getlength(text) / MEASURE_SCALE


def assert_single_line() -> None:
    """Fail the build if any list item would wrap onto a second line."""
    inches = int(TEXT_COL_W) - int(LIST_INSET_L) - int(LIST_INSET_R) - int(BULLET_INDENT)
    available_pt = inches / 914400 * 72 - FIT_SAFETY_PT
    problems = []
    for paper in content.PAPERS:
        for _style, label, text in paper["items"]:
            width = text_width_pt(f"{label}: ", SANS, True, LABEL_PT)
            width += text_width_pt(text, SANS, False, BODY_PT)
            if width > available_pt:
                problems.append(f"  paper {paper['number']:>2} {label}: {width:.1f}pt > {available_pt:.1f}pt -- {text}")
    if problems:
        raise SystemExit("single-line fit check failed:\n" + "\n".join(problems))


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill: RGBColor | None):
    """A borderless rectangle used for rules, bands and panels."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
    return shape


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return box, frame


def style_run(run, *, size, bold=False, italic=False, color=INK, typeface=SANS):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = typeface


def simple_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    typeface=SANS,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.0,
):
    box, frame = add_textbox(slide, left, top, width, height)
    frame.vertical_anchor = anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = spacing
    run = paragraph.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, italic=italic, color=color, typeface=typeface)
    return box


def set_bullet(paragraph, style: str, first_numbered: bool) -> None:
    """Attach a square bullet or an arabic auto-number to ``paragraph``."""
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("marL", str(int(BULLET_INDENT)))
    p_pr.set("indent", str(-int(BULLET_INDENT)))
    p_pr.append(parse_xml(f'<a:buClr {nsdecls("a")}><a:srgbClr val="407CB9"/></a:buClr>'))
    if style == "number":
        p_pr.append(parse_xml(f'<a:buFont {nsdecls("a")} typeface="{SANS}"/>'))
        start = ' startAt="1"' if first_numbered else ""
        p_pr.append(parse_xml(f'<a:buAutoNum {nsdecls("a")} type="arabicPeriod"{start}/>'))
    else:
        p_pr.append(parse_xml(f'<a:buFont {nsdecls("a")} typeface="{SANS}"/>'))
        p_pr.append(parse_xml(f'<a:buChar {nsdecls("a")} char="{BULLET_CHAR}"/>'))


def add_hyperlink_run(paragraph, text: str, address: str, *, size, color=BLUE):
    run = paragraph.add_run()
    run.text = text
    style_run(run, size=size, color=color)
    run.hyperlink.address = address
    return run


# ---------------------------------------------------------------------------
# Slide furniture shared by every slide
# ---------------------------------------------------------------------------
def add_chrome(slide, page_number: int, *, header_right: str = "", show_header: bool = True):
    if show_header:
        simple_text(
            slide,
            MARGIN,
            HEADER_Y,
            Inches(8.6),
            Inches(0.24),
            content.DECK_TITLE,
            size=9,
            color=GRAY,
        )
        if header_right:
            simple_text(
                slide,
                Inches(8.9),
                HEADER_Y,
                Emu(int(CONTENT_RIGHT) - int(Inches(8.9))),
                Inches(0.24),
                header_right,
                size=9,
                color=BLUE,
                align=PP_ALIGN.RIGHT,
            )
        add_rect(
            slide,
            MARGIN,
            RULE_Y,
            Emu(int(CONTENT_RIGHT) - int(MARGIN)),
            Inches(0.035),
            NAVY,
        )

    add_rect(
        slide,
        MARGIN,
        FOOTER_RULE_Y,
        Emu(int(CONTENT_RIGHT) - int(MARGIN)),
        Inches(0.012),
        SILVER,
    )
    simple_text(
        slide,
        MARGIN,
        FOOTER_TEXT_Y,
        Inches(5.2),
        FOOTER_H,
        content.AUTHOR_SHORT,
        size=9,
        color=GRAY,
    )
    simple_text(
        slide,
        Inches(5.4),
        FOOTER_TEXT_Y,
        Inches(4.0),
        FOOTER_H,
        content.PRESENTED_DATE,
        size=9,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    simple_text(
        slide,
        Inches(11.6),
        FOOTER_TEXT_Y,
        Emu(int(CONTENT_RIGHT) - int(Inches(11.6))),
        FOOTER_H,
        str(page_number),
        size=10,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def build_title_slide(presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, Emu(0), Emu(0), Inches(0.30), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.30), Emu(0), Inches(0.07), SLIDE_H, BLUE)

    left = Inches(1.05)
    width = Emu(int(CONTENT_RIGHT) - int(left))

    simple_text(
        slide,
        left,
        Inches(0.86),
        width,
        Inches(0.3),
        "CHEMICALQDEVICE  |  INDEPENDENT RESEARCH",
        size=11,
        color=BLUE,
        bold=True,
    )

    top = Inches(1.34)
    for line in content.DECK_TITLE_LINES:
        simple_text(
            slide,
            left,
            top,
            width,
            Inches(0.62),
            line,
            size=31,
            color=NAVY,
            bold=True,
            typeface=SERIF,
        )
        top = Emu(int(top) + int(Inches(0.58)))

    add_rect(slide, left, Inches(2.66), Inches(2.1), Inches(0.045), BLUE)

    simple_text(
        slide,
        left,
        Inches(2.88),
        width,
        Inches(0.34),
        content.DECK_SUBTITLE,
        size=15,
        color=GRAY,
        italic=True,
        typeface=SERIF,
    )

    stat_top = Inches(3.46)
    stat_h = Inches(1.16)
    gap = Inches(0.22)
    stat_w = Emu(int((int(width) - 3 * int(gap)) / 4))
    for index, (value, label, note) in enumerate(content.TITLE_STATS):
        x = Emu(int(left) + index * (int(stat_w) + int(gap)))
        add_rect(slide, x, stat_top, stat_w, stat_h, PALE)
        add_rect(slide, x, stat_top, Inches(0.05), stat_h, NAVY)
        simple_text(
            slide,
            Emu(int(x) + int(Inches(0.22))),
            Emu(int(stat_top) + int(Inches(0.13))),
            Emu(int(stat_w) - int(Inches(0.34))),
            Inches(0.42),
            value,
            size=24,
            color=NAVY,
            bold=True,
            typeface=SERIF,
        )
        simple_text(
            slide,
            Emu(int(x) + int(Inches(0.22))),
            Emu(int(stat_top) + int(Inches(0.58))),
            Emu(int(stat_w) - int(Inches(0.34))),
            Inches(0.26),
            label,
            size=10.5,
            color=INK,
            bold=True,
        )
        simple_text(
            slide,
            Emu(int(x) + int(Inches(0.22))),
            Emu(int(stat_top) + int(Inches(0.82))),
            Emu(int(stat_w) - int(Inches(0.34))),
            Inches(0.26),
            note,
            size=9.5,
            color=GRAY,
        )

    simple_text(slide, left, Inches(4.92), width, Inches(0.32), content.AUTHOR, size=15, color=INK, bold=True)
    simple_text(slide, left, Inches(5.24), width, Inches(0.28), content.CONTACT, size=11, color=GRAY)

    add_rect(slide, left, Inches(5.66), Emu(int(width)), Inches(0.012), SILVER)

    simple_text(
        slide,
        left,
        Inches(5.82),
        width,
        Inches(0.30),
        content.VENUE,
        size=12.5,
        color=NAVY,
        bold=True,
    )
    simple_text(
        slide,
        left,
        Inches(6.12),
        width,
        Inches(0.30),
        content.PRESENTED_DATE,
        size=12.5,
        color=GRAY,
    )

    box, frame = add_textbox(slide, left, Inches(6.46), width, Inches(0.36))
    paragraph = frame.paragraphs[0]
    paragraph.line_spacing = 1.08
    run = paragraph.add_run()
    run.text = content.DISCLAIMER
    style_run(run, size=8.5, italic=True, color=GRAY)

    add_chrome(slide, 1, show_header=False)


def place_cover(slide, image_path: Path, side: str):
    with Image.open(image_path) as image:
        width_px, height_px = image.size
    aspect = width_px / height_px

    box_w = int(IMAGE_BOX_W)
    box_h = int(IMAGE_BOX_H)
    draw_h = box_h
    draw_w = int(draw_h * aspect)
    if draw_w > box_w:
        draw_w = box_w
        draw_h = int(draw_w / aspect)

    region_top = int(BODY_TOP)
    region_h = int(BODY_BOTTOM) - int(BODY_TOP)
    top = region_top + (region_h - draw_h) // 2

    if side == "left":
        box_left = int(MARGIN)
    else:
        box_left = int(CONTENT_RIGHT) - box_w
    left = box_left + (box_w - draw_w) // 2

    picture = slide.shapes.add_picture(str(image_path), Emu(left), Emu(top), Emu(draw_w), Emu(draw_h))
    picture.line.color.rgb = SILVER
    picture.line.width = Pt(0.75)
    return picture


def build_paper_slide(presentation, paper: dict) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)

    number = paper["number"]
    side = content.side_for(number)
    add_chrome(
        slide,
        number + 1,
        header_right=f"Paper {number:02d} of 20  |  {paper['doi']}",
    )

    simple_text(
        slide,
        MARGIN,
        HEADLINE_Y,
        Emu(int(CONTENT_RIGHT) - int(MARGIN)),
        HEADLINE_H,
        f"{number:02d}.  {paper['headline']}",
        size=21,
        color=NAVY,
        bold=True,
        typeface=SERIF,
    )
    add_rect(slide, MARGIN, Inches(1.04), Inches(1.5), Inches(0.028), BLUE)

    place_cover(slide, COVER_DIR / paper["cover"], side)

    if side == "left":
        text_left = Emu(int(MARGIN) + int(IMAGE_BOX_W) + int(COLUMN_GAP))
    else:
        text_left = MARGIN

    items = list(paper["items"])
    synthesis = [item for item in items if item[1].startswith("LLM ")]
    evidence = [item for item in items if not item[1].startswith("LLM ")]

    panel_h = Inches(0.44 + 0.34 * len(synthesis))
    panel_top = Emu(int(BODY_BOTTOM) - int(panel_h))
    add_rect(slide, text_left, panel_top, TEXT_COL_W, panel_h, PALE)
    add_rect(slide, text_left, panel_top, Inches(0.05), panel_h, NAVY)

    evidence_h = Emu(int(panel_top) - int(BODY_TOP) - int(Inches(0.22)))
    box, frame = add_textbox(slide, text_left, BODY_TOP, TEXT_COL_W, evidence_h)
    frame.margin_left = LIST_INSET_L
    frame.margin_right = LIST_INSET_R
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    seen_number = False
    for index, (style, label, text) in enumerate(evidence):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.line_spacing = 1.02
        paragraph.space_after = Pt(28)
        set_bullet(paragraph, style, first_numbered=(style == "number" and not seen_number))
        if style == "number":
            seen_number = True
        label_run = paragraph.add_run()
        label_run.text = f"{label}: "
        style_run(label_run, size=LABEL_PT, bold=True, color=NAVY)
        body_run = paragraph.add_run()
        body_run.text = text
        style_run(body_run, size=BODY_PT, color=INK)

    panel_box, panel_frame = add_textbox(
        slide,
        text_left,
        Emu(int(panel_top) + int(Inches(0.16))),
        TEXT_COL_W,
        Emu(int(panel_h) - int(Inches(0.32))),
    )
    panel_frame.margin_left = LIST_INSET_L
    panel_frame.margin_right = LIST_INSET_R
    panel_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for index, (style, label, text) in enumerate(synthesis):
        paragraph = panel_frame.paragraphs[0] if index == 0 else panel_frame.add_paragraph()
        paragraph.line_spacing = 1.02
        paragraph.space_after = Pt(6)
        set_bullet(paragraph, style, first_numbered=False)
        label_run = paragraph.add_run()
        label_run.text = f"{label}: "
        style_run(label_run, size=LABEL_PT, bold=True, color=BLUE)
        body_run = paragraph.add_run()
        body_run.text = text
        style_run(body_run, size=BODY_PT, color=INK)


def build_reference_slide(presentation, papers, page_number: int, part: int, total: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
    add_chrome(
        slide,
        page_number,
        header_right=f"References  |  part {part} of {total}",
    )

    simple_text(
        slide,
        MARGIN,
        HEADLINE_Y,
        Emu(int(CONTENT_RIGHT) - int(MARGIN)),
        HEADLINE_H,
        f"References ({papers[0]['number']} to {papers[-1]['number']} of 20)",
        size=21,
        color=NAVY,
        bold=True,
        typeface=SERIF,
    )
    add_rect(slide, MARGIN, Inches(1.04), Inches(1.5), Inches(0.028), BLUE)

    box, frame = add_textbox(
        slide,
        MARGIN,
        BODY_TOP,
        Emu(int(CONTENT_RIGHT) - int(MARGIN)),
        Emu(int(BODY_BOTTOM) - int(BODY_TOP) - int(Inches(0.36))),
    )
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, paper in enumerate(papers):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.line_spacing = 1.04
        paragraph.space_after = Pt(15)
        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("marL", str(int(Inches(0.42))))
        p_pr.set("indent", str(-int(Inches(0.42))))

        number_run = paragraph.add_run()
        number_run.text = f"{paper['number']}.\t"
        style_run(number_run, size=11, bold=True, color=NAVY)

        year = paper["date"].rsplit(" ", 1)[-1]
        author_run = paragraph.add_run()
        author_run.text = f"Kawchak, K. ({year}). "
        style_run(author_run, size=11, color=INK)

        title_run = paragraph.add_run()
        title_run.text = f"{paper['reference_title']}. "
        style_run(title_run, size=11, italic=True, color=INK)

        source_run = paragraph.add_run()
        source_run.text = "Zenodo. "
        style_run(source_run, size=11, color=GRAY)

        doi = paper["doi"]
        add_hyperlink_run(paragraph, f"https://doi.org/{doi}", f"https://doi.org/{doi}", size=11)

    simple_text(
        slide,
        MARGIN,
        Inches(6.38),
        Emu(int(CONTENT_RIGHT) - int(MARGIN)),
        Inches(0.30),
        content.CLOSING_NOTE,
        size=9.5,
        color=GRAY,
        italic=True,
    )


def build() -> Path:
    assert_single_line()

    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H

    build_title_slide(presentation)
    for paper in content.PAPERS:
        build_paper_slide(presentation, paper)

    half = len(content.PAPERS) // 2
    build_reference_slide(presentation, content.PAPERS[:half], len(content.PAPERS) + 2, 1, 2)
    build_reference_slide(presentation, content.PAPERS[half:], len(content.PAPERS) + 3, 2, 2)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


def main() -> int:
    path = build()
    print(f"wrote {path.relative_to(REPO_ROOT)} ({path.stat().st_size:,} bytes)")
    print(f"{len(content.PAPERS) + 3} slides: 1 title, {len(content.PAPERS)} papers, 2 references")
    return 0


if __name__ == "__main__":
    sys.exit(main())
